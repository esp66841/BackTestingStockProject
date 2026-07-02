"""
Author: Elias Powell
Date: 2026-04-25

About:
This is a Python project for stock strategy design and backtesting on the U.S. stock market.

Strategy:
The strategy buys a stock after it drops 5% or more over the past 5 trading days.
To avoid look-ahead bias, the buy signal is created after the close and the trade is executed on the next trading day.
Instead of holding for exactly 5 trading days, the position is held until the stock rises 3% or more from the entry price.
If the 3% target is not reached by the end of the backtest, the position is sold on the final trading day.

Data Source: Yahoo Finance through the yfinance Python package.
Selected Stocks: 20 different S&P 500 stocks.
Trading Period: 2012-01-01 to 2025-12-31.
Initial Capital: $100,000 per stock backtest.
Position Sizing: Invest 100% of available cash when a valid buy signal occurs.
Transaction Costs: Ignored for simplicity.
Constraint: The portfolio balance cannot become negative because the strategy only buys when cash is available.
Output: Summary-only Excel file and PowerPoint-ready visuals based on the summary results.
"""

# Importing necessary libraries
import math
import os
import tempfile
from datetime import datetime

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import yfinance as yf


tickers = [
    "ANET", "TMO", "TMUS", "QCOM", "BA",
    "VZ", "TJX", "SCHW", "UNP", "WELL",
    "UBER", "BX", "CB", "PLD", "SPGI",
    "COF", "LMT", "SYK", "SBUX", "FTNT",
]

start_date = "2012-01-01"
end_date = "2025-12-31"
download_end_date = "2026-01-01"

initial_capital = 100000.0
drop_threshold = -0.05
profit_target = 0.03

excel_output_file = "backtesting_results_20_stocks_summary_only.xlsx"
visuals_folder = "ppt_visuals_20_stocks"
ppt_output_file = "stock_backtesting_summary_visuals_20_stocks.pptx"

yfinance_cache_dir = os.path.join(
    tempfile.gettempdir(), "python_stock_backtest_yfinance_cache"
)
os.makedirs(yfinance_cache_dir, exist_ok=True)
yf.set_tz_cache_location(yfinance_cache_dir)


def download_stock_data(ticker, start_date, end_date):
    """Download stock data using yfinance."""
    data = yf.download(
        ticker,
        start=start_date,
        end=end_date,
        auto_adjust=False,
        progress=False,
    )

    if isinstance(data.columns, pd.MultiIndex):
        if ticker in data.columns.get_level_values(-1):
            data = data.xs(ticker, axis=1, level=-1)
        else:
            data.columns = data.columns.get_level_values(0)

    data = data.reset_index()
    data["Ticker"] = ticker

    return data


def clean_stock_data(data):
    """Clean the data and calculate daily and 5-day returns from Close prices."""
    data = data.copy()
    data["Date"] = pd.to_datetime(data["Date"])
    data = data.dropna(subset=["Close"])
    data = data.sort_values("Date")
    data = data[(data["Date"] >= start_date) & (data["Date"] <= end_date)]

    data["Daily_Return"] = data["Close"].pct_change()
    data["Five_Day_Return"] = data["Close"].pct_change(periods=5)
    data = data.dropna(subset=["Five_Day_Return"]).reset_index(drop=True)

    return data


def get_execution_price(row):
    """Use Open for trade execution when it is available, otherwise use Close."""
    if "Open" in row.index and pd.notna(row["Open"]) and row["Open"] > 0:
        return float(row["Open"])

    return float(row["Close"])


def calculate_max_drawdown(value_series):
    running_peak = value_series.cummax()
    drawdown = value_series / running_peak - 1
    return drawdown.min()


def calculate_sharpe_ratio(daily_returns):
    daily_returns = daily_returns.dropna()

    if daily_returns.empty or daily_returns.std() == 0:
        return np.nan

    return daily_returns.mean() / daily_returns.std() * math.sqrt(252)


def backtest_drop_strategy(
    data,
    initial_capital=100000.0,
    drop_threshold=-0.05,
    profit_target=0.03,
):
    """
    Backtest the 5-day drop strategy with a 3% profit target.

    The buy signal is observed at the close and shifted forward one trading day.
    The 3% gain is also observed at the close, then the sell executes on the
    next available trading day Open when possible.
    """
    data = data.copy().reset_index(drop=True)
    data["Buy_Signal"] = np.where(data["Five_Day_Return"] <= drop_threshold, 1, 0)
    data["Executable_Buy_Signal"] = data["Buy_Signal"].shift(1).fillna(0).astype(int)

    data["Cash"] = float(initial_capital)
    data["Shares"] = 0.0
    data["Holdings"] = 0.0
    data["Portfolio_Value"] = float(initial_capital)
    data["Trade"] = ""
    data["Entry_Price"] = np.nan
    data["Exit_Price"] = np.nan
    data["Return_From_Entry"] = np.nan
    data["Days_Held"] = np.nan

    cash = float(initial_capital)
    shares = 0.0
    entry_price = np.nan
    entry_index = None
    pending_sell = False
    closed_trades = []

    for i in range(len(data)):
        row = data.iloc[i]
        daily_trades = []
        exit_price = np.nan
        current_return_from_entry = np.nan
        current_days_held = np.nan

        if pending_sell and shares > 0:
            exit_price = get_execution_price(row)
            cash = shares * exit_price
            exit_return = (exit_price / entry_price) - 1
            current_days_held = i - entry_index

            daily_trades.append("Sell")
            closed_trades.append(
                {
                    "Return_From_Entry": exit_return,
                    "Days_Held": current_days_held,
                    "Trade": "Sell",
                }
            )

            shares = 0.0
            entry_price = np.nan
            entry_index = None
            pending_sell = False

        if shares == 0 and cash > 0 and int(row["Executable_Buy_Signal"]) == 1:
            buy_price = get_execution_price(row)

            if buy_price > 0:
                shares = cash / buy_price
                cash = 0.0
                entry_price = buy_price
                entry_index = i
                daily_trades.append("Buy")

        if shares > 0:
            current_return_from_entry = (row["Close"] / entry_price) - 1
            current_days_held = i - entry_index

            # The sell signal is observed at the close. The actual sell happens
            # on the next available trading day using Open when possible.
            if current_return_from_entry >= profit_target and i < len(data) - 1:
                pending_sell = True

        if i == len(data) - 1 and shares > 0:
            exit_price = float(row["Close"])
            cash = shares * exit_price
            exit_return = (exit_price / entry_price) - 1
            current_days_held = i - entry_index

            daily_trades.append("Final Sell")
            closed_trades.append(
                {
                    "Return_From_Entry": exit_return,
                    "Days_Held": current_days_held,
                    "Trade": "Final Sell",
                }
            )

            shares = 0.0
            entry_price = np.nan
            entry_index = None
            pending_sell = False
            current_return_from_entry = exit_return

        holdings = shares * float(row["Close"])
        portfolio_value = cash + holdings

        data.loc[i, "Cash"] = float(cash)
        data.loc[i, "Shares"] = float(shares)
        data.loc[i, "Holdings"] = float(holdings)
        data.loc[i, "Portfolio_Value"] = float(portfolio_value)
        data.loc[i, "Trade"] = " / ".join(daily_trades)
        data.loc[i, "Entry_Price"] = entry_price if shares > 0 else np.nan
        data.loc[i, "Exit_Price"] = exit_price
        data.loc[i, "Return_From_Entry"] = current_return_from_entry
        data.loc[i, "Days_Held"] = current_days_held

    data["Strategy_Return"] = data["Portfolio_Value"].pct_change().fillna(0)
    data["Strategy_Cumulative_Return"] = (1 + data["Strategy_Return"]).cumprod() - 1

    return data, pd.DataFrame(closed_trades)


def backtest_buy_and_hold(data, initial_capital=100000.0):
    """Create a buy-and-hold benchmark for comparison."""
    data = data.copy()
    first_price = data["Close"].iloc[0]
    shares = initial_capital / first_price

    data["Buy_Hold_Value"] = shares * data["Close"]
    data["Buy_Hold_Return"] = data["Buy_Hold_Value"].pct_change().fillna(0)

    return data


def calculate_metrics(data, closed_trade_log):
    """Calculate strategy and benchmark performance metrics."""
    strategy_start = data["Portfolio_Value"].iloc[0]
    strategy_end = data["Portfolio_Value"].iloc[-1]

    benchmark_start = data["Buy_Hold_Value"].iloc[0]
    benchmark_end = data["Buy_Hold_Value"].iloc[-1]

    strategy_total_return = (strategy_end / strategy_start) - 1
    benchmark_total_return = (benchmark_end / benchmark_start) - 1

    strategy_annual_return = (1 + strategy_total_return) ** (252 / len(data)) - 1
    benchmark_annual_return = (1 + benchmark_total_return) ** (252 / len(data)) - 1

    strategy_volatility = data["Strategy_Return"].std() * np.sqrt(252)
    benchmark_volatility = data["Buy_Hold_Return"].std() * np.sqrt(252)

    if closed_trade_log.empty:
        number_of_trades = 0
        average_days_held = np.nan
        winning_trades = 0
        losing_trades = 0
        win_rate = np.nan
    else:
        number_of_trades = len(closed_trade_log)
        average_days_held = closed_trade_log["Days_Held"].mean()
        winning_trades = int((closed_trade_log["Return_From_Entry"] > 0).sum())
        losing_trades = int((closed_trade_log["Return_From_Entry"] <= 0).sum())
        win_rate = winning_trades / number_of_trades

    return {
        "Final Strategy Value": strategy_end,
        "Final Buy Hold Value": benchmark_end,
        "Strategy Total Return": strategy_total_return,
        "Buy Hold Total Return": benchmark_total_return,
        "Strategy Annual Return": strategy_annual_return,
        "Buy Hold Annual Return": benchmark_annual_return,
        "Strategy Volatility": strategy_volatility,
        "Buy Hold Volatility": benchmark_volatility,
        "Strategy Sharpe Ratio": calculate_sharpe_ratio(data["Strategy_Return"]),
        "Buy Hold Sharpe Ratio": calculate_sharpe_ratio(data["Buy_Hold_Return"]),
        "Strategy Max Drawdown": calculate_max_drawdown(data["Portfolio_Value"]),
        "Buy Hold Max Drawdown": calculate_max_drawdown(data["Buy_Hold_Value"]),
        "Number of Trades": number_of_trades,
        "Average Days Held": average_days_held,
        "Winning Trades": winning_trades,
        "Losing Trades": losing_trades,
        "Win Rate": win_rate,
    }


def export_summary_results(results_df):
    """Export only the Summary Results sheet."""
    output_file = excel_output_file

    try:
        writer = pd.ExcelWriter(output_file, engine="openpyxl")
    except PermissionError:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = f"backtesting_results_20_stocks_summary_only_{timestamp}.xlsx"
        print(
            f"{excel_output_file} is currently locked. "
            f"Saving Excel results to {output_file} instead."
        )
        writer = pd.ExcelWriter(output_file, engine="openpyxl")

    with writer:
        results_df.to_excel(writer, sheet_name="Summary Results", index=False)

    return output_file


def format_percent_axis(ax):
    ax.xaxis.set_major_formatter(lambda x, pos: f"{x:.0%}")


def save_strategy_vs_buy_hold_total_return(results_df):
    chart_data = results_df.sort_values("Strategy Total Return")
    y_positions = np.arange(len(chart_data))
    bar_height = 0.38

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.barh(
        y_positions - bar_height / 2,
        chart_data["Strategy Total Return"],
        height=bar_height,
        label="Strategy",
    )
    ax.barh(
        y_positions + bar_height / 2,
        chart_data["Buy Hold Total Return"],
        height=bar_height,
        label="Buy and Hold",
    )
    ax.set_yticks(y_positions)
    ax.set_yticklabels(chart_data["Ticker"])
    ax.set_title("Strategy vs Buy-and-Hold Total Return")
    ax.set_xlabel("Total Return")
    format_percent_axis(ax)
    ax.legend()
    ax.grid(True, axis="x", alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "strategy_vs_buy_hold_total_return.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_final_portfolio_value_comparison(results_df):
    chart_data = results_df.sort_values("Final Strategy Value")
    y_positions = np.arange(len(chart_data))
    bar_height = 0.38

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.barh(
        y_positions - bar_height / 2,
        chart_data["Final Strategy Value"],
        height=bar_height,
        label="Strategy",
    )
    ax.barh(
        y_positions + bar_height / 2,
        chart_data["Final Buy Hold Value"],
        height=bar_height,
        label="Buy and Hold",
    )
    ax.set_yticks(y_positions)
    ax.set_yticklabels(chart_data["Ticker"])
    ax.set_title("Final Portfolio Value Comparison")
    ax.set_xlabel("Portfolio Value ($)")
    ax.xaxis.set_major_formatter(lambda x, pos: f"${x / 1000:.0f}K")
    ax.legend()
    ax.grid(True, axis="x", alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "final_portfolio_value_comparison.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_strategy_max_drawdown(results_df):
    chart_data = results_df.sort_values("Strategy Max Drawdown")

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.bar(chart_data["Ticker"], chart_data["Strategy Max Drawdown"])
    ax.set_title("Strategy Max Drawdown by Ticker")
    ax.set_xlabel("Ticker")
    ax.set_ylabel("Max Drawdown")
    ax.yaxis.set_major_formatter(lambda y, pos: f"{y:.0%}")
    ax.tick_params(axis="x", rotation=45)
    ax.grid(True, axis="y", alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "strategy_max_drawdown.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_strategy_win_rate(results_df):
    chart_data = results_df.sort_values("Win Rate")

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.bar(chart_data["Ticker"], chart_data["Win Rate"])
    ax.set_title("Strategy Win Rate by Ticker")
    ax.set_xlabel("Ticker")
    ax.set_ylabel("Win Rate")
    ax.yaxis.set_major_formatter(lambda y, pos: f"{y:.0%}")
    ax.set_ylim(0, 1.05)
    ax.tick_params(axis="x", rotation=45)
    ax.grid(True, axis="y", alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "strategy_win_rate.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_average_days_held(results_df):
    chart_data = results_df.sort_values("Average Days Held", ascending=False)

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.bar(chart_data["Ticker"], chart_data["Average Days Held"])
    ax.set_title("Average Days Held by Ticker")
    ax.set_xlabel("Ticker")
    ax.set_ylabel("Average Days Held")
    ax.tick_params(axis="x", rotation=45)
    ax.grid(True, axis="y", alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "average_days_held.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_risk_return_scatter(results_df):
    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.scatter(results_df["Strategy Volatility"], results_df["Strategy Annual Return"], s=70)

    for _, row in results_df.iterrows():
        ax.annotate(
            row["Ticker"],
            (row["Strategy Volatility"], row["Strategy Annual Return"]),
            xytext=(5, 5),
            textcoords="offset points",
            fontsize=9,
        )

    ax.set_title("Strategy Risk vs Return")
    ax.set_xlabel("Strategy Volatility")
    ax.set_ylabel("Strategy Annual Return")
    ax.xaxis.set_major_formatter(lambda x, pos: f"{x:.0%}")
    ax.yaxis.set_major_formatter(lambda y, pos: f"{y:.0%}")
    ax.grid(True, alpha=0.3)
    fig.tight_layout()
    path = os.path.join(visuals_folder, "risk_return_scatter.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def save_summary_metrics_table(results_df):
    table_data = results_df.sort_values(
        "Strategy Total Return", ascending=False
    ).head(10)
    display_data = pd.DataFrame(
        {
            "Ticker": table_data["Ticker"],
            "Strategy Total Return": table_data["Strategy Total Return"].map(
                lambda x: f"{x:.1%}"
            ),
            "Buy Hold Total Return": table_data["Buy Hold Total Return"].map(
                lambda x: f"{x:.1%}"
            ),
            "Strategy Max Drawdown": table_data["Strategy Max Drawdown"].map(
                lambda x: f"{x:.1%}"
            ),
            "Win Rate": table_data["Win Rate"].map(
                lambda x: "" if pd.isna(x) else f"{x:.1%}"
            ),
            "Average Days Held": table_data["Average Days Held"].map(
                lambda x: "" if pd.isna(x) else f"{x:.1f}"
            ),
        }
    )

    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.axis("off")
    ax.set_title("Top 10 Stocks by Strategy Total Return", fontsize=16, pad=20)
    table = ax.table(
        cellText=display_data.values,
        colLabels=display_data.columns,
        loc="center",
        cellLoc="center",
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1, 1.5)

    for (row, col), cell in table.get_celld().items():
        if row == 0:
            cell.set_text_props(weight="bold")
            cell.set_facecolor("#E8EEF7")
        elif row % 2 == 0:
            cell.set_facecolor("#F7F9FC")

    fig.tight_layout()
    path = os.path.join(visuals_folder, "summary_metrics_table.png")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return path


def create_summary_visuals(results_df):
    """Create PowerPoint-ready visuals from the summary results only."""
    os.makedirs(visuals_folder, exist_ok=True)

    chart_paths = [
        save_strategy_vs_buy_hold_total_return(results_df),
        save_final_portfolio_value_comparison(results_df),
        save_strategy_max_drawdown(results_df),
        save_strategy_win_rate(results_df),
        save_average_days_held(results_df),
        save_risk_return_scatter(results_df),
        save_summary_metrics_table(results_df),
    ]

    return chart_paths


def create_optional_powerpoint(chart_paths):
    """Create a PowerPoint if python-pptx is installed."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("python-pptx is not installed; skipping PowerPoint creation.")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Stock Backtesting Summary Visuals"
    title_slide.placeholders[1].text = (
        "20-stock summary for the 5-day drop and 3% profit target strategy"
    )

    for chart_path in chart_paths:
        title = os.path.splitext(os.path.basename(chart_path))[0].replace("_", " ").title()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(chart_path, Inches(0.65), Inches(1.15), width=Inches(12.0))

    prs.save(ppt_output_file)
    return ppt_output_file


def main():
    all_results = []

    for ticker in tickers:
        print(f"Running backtest for {ticker}...")

        data = download_stock_data(ticker, start_date, download_end_date)
        data = clean_stock_data(data)

        if data.empty:
            print(f"No usable data found for {ticker}; skipping.")
            continue

        data, closed_trade_log = backtest_drop_strategy(
            data,
            initial_capital=initial_capital,
            drop_threshold=drop_threshold,
            profit_target=profit_target,
        )
        data = backtest_buy_and_hold(data, initial_capital)

        metrics = calculate_metrics(data, closed_trade_log)
        metrics["Ticker"] = ticker
        all_results.append(metrics)

    if not all_results:
        print("No backtests were completed because no ticker data was loaded.")
        return

    results_df = pd.DataFrame(all_results)
    ordered_columns = [
        "Ticker",
        "Final Strategy Value",
        "Final Buy Hold Value",
        "Strategy Total Return",
        "Buy Hold Total Return",
        "Strategy Annual Return",
        "Buy Hold Annual Return",
        "Strategy Volatility",
        "Buy Hold Volatility",
        "Strategy Sharpe Ratio",
        "Buy Hold Sharpe Ratio",
        "Strategy Max Drawdown",
        "Buy Hold Max Drawdown",
        "Number of Trades",
        "Average Days Held",
        "Winning Trades",
        "Losing Trades",
        "Win Rate",
    ]
    results_df = results_df[ordered_columns]

    print(results_df)
    saved_excel_file = export_summary_results(results_df)
    chart_paths = create_summary_visuals(results_df)
    saved_ppt_file = create_optional_powerpoint(chart_paths)

    print(f"\nExcel results exported to {saved_excel_file}")
    print(f"PNG visuals saved in {visuals_folder}")

    if saved_ppt_file:
        print(f"PowerPoint exported to {saved_ppt_file}")


if __name__ == "__main__":
    main()
